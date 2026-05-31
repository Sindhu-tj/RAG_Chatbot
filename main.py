import streamlit as st
import tempfile
import os
import json
import re
import pandas as pd
from PIL import Image
from docx import Document
from pptx import Presentation
import easyocr
from langchain.schema import Document as LangDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from transformers import pipeline

# =========================================================
# PAGE CONFIG
# =========================================================
st.set_page_config(page_title="Universal AI File Chatbot", page_icon="🤖", layout="wide")
st.markdown("""
<style>
.main { background-color: #0E1117; }
.stDataFrame { border-radius: 12px; }
.answer-card {
    background: #1a1d2e;
    border-left: 4px solid #4f8ef7;
    border-radius: 10px;
    padding: 18px 22px;
    margin-top: 8px;
}
</style>
""", unsafe_allow_html=True)

# =========================================================
# SESSION STATE
# =========================================================
DEFAULTS = {"last_file": None, "df": None, "text": "", "file_type": "", "vectorstore": None}
for k, v in DEFAULTS.items():
    if k not in st.session_state:
        st.session_state[k] = v

TABULAR   = ["csv", "xlsx", "xls"]
IMAGE_EXT = ["png", "jpg", "jpeg"]

# =========================================================
# CACHED RESOURCES
# =========================================================
@st.cache_resource
def load_ocr():
    return easyocr.Reader(['en'], gpu=False)

@st.cache_resource
def load_model():
    # Try multiple task names for compatibility across transformers versions
    task_options = [
        ("text2text-generation", {"model": "google/flan-t5-large", "max_new_tokens": 256}),
        ("text-generation",      {"model": "gpt2", "max_new_tokens": 256}),
        ("summarization",        {"model": "facebook/bart-large-cnn", "max_length": 256}),
    ]
    for task, kwargs in task_options:
        try:
            return pipeline(task, **kwargs), task
        except Exception:
            continue
    raise RuntimeError("Could not load any text generation model. Run: pip install transformers --upgrade")

# =========================================================
# INTENT DETECTOR
# =========================================================
def detect_intent(question: str) -> str:
    q = question.lower().strip()

    LIST_KW    = ["list", "show all", "give all", "display all", "all the",
                  "fetch all", "get all", "enumerate", "what are all",
                  "names of all", "show me all", "give me all", "all names",
                  "show the list", "give the list", "what are the"]
    COUNT_KW   = ["how many", "count", "total number", "number of", "how much",
                  "total count", "quantity"]
    TABLE_KW   = ["find", "search for", "filter", "which records", "rows where",
                  "entries where", "show records", "look for"]
    SUMMARY_KW = ["summary", "overview", "describe", "about this file",
                  "what is this", "tell me about", "what does this file contain",
                  "information about the file", "what's in this"]
    COL_KW     = ["column", "columns", "fields", "headers", "attributes"]

    if any(kw in q for kw in LIST_KW):    return "list"
    if any(kw in q for kw in COUNT_KW):   return "count"
    if any(kw in q for kw in SUMMARY_KW): return "summary"
    if any(kw in q for kw in COL_KW):     return "columns"
    if any(kw in q for kw in TABLE_KW):   return "table"
    return "factual"

# =========================================================
# COLUMN MATCHER — finds which df column the user means
# =========================================================
def find_column(question: str, df: pd.DataFrame):
    q = question.lower()
    col_map = {c.lower(): c for c in df.columns}

    # Exact substring match
    for cl, cr in col_map.items():
        if cl in q:
            return cr

    # Word-level match
    for cl, cr in col_map.items():
        words = re.split(r'[\s_\-]+', cl)
        if any(w in q for w in words if len(w) > 2):
            return cr

    # Single column → return it
    if len(df.columns) == 1:
        return df.columns[0]

    return None

# =========================================================
# FORMATTERS — always return clean markdown strings
# =========================================================
def fmt_list(values, label="Results"):
    items = sorted({str(v).strip() for v in values if str(v).strip()})
    if not items:
        return "No items found."
    header = f"**{label}** — {len(items)} items\n"
    body   = "\n".join(f"- {item}" for item in items)
    return header + "\n" + body

def fmt_count(value, label=""):
    return f"**{label + ': ' if label else ''}{value}**"

def fmt_columns(df):
    lines = ["**File columns:**\n"] + [f"- {c}" for c in df.columns]
    return "\n".join(lines)

def fmt_summary(df):
    lines = [
        f"**Rows:** {df.shape[0]}",
        f"**Columns:** {df.shape[1]}",
        "\n**Column names:**",
    ] + [f"- {c}" for c in df.columns]
    return "\n".join(lines)

# =========================================================
# TEXT EXTRACTION  (all file types)
# =========================================================
def extract_text(uploaded_file, file_type: str) -> str:
    text = ""
    try:
        if file_type == "pdf":
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(uploaded_file.read())
                path = tmp.name
            docs = PyPDFLoader(path).load()
            text = "\n".join(d.page_content for d in docs)
            os.unlink(path)

        elif file_type == "docx":
            doc  = Document(uploaded_file)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif file_type in ["txt", "md", "xml"]:
            text = uploaded_file.read().decode("utf-8", errors="ignore")

        elif file_type == "json":
            text = uploaded_file.read().decode("utf-8", errors="ignore")

        elif file_type == "csv":
            df = pd.read_csv(uploaded_file).fillna("")
            st.session_state.df = df
            text = df.to_string(index=False)

        elif file_type in ["xlsx", "xls"]:
            df = pd.read_excel(uploaded_file).fillna("")
            st.session_state.df = df
            text = df.to_string(index=False)

        elif file_type == "pptx":
            prs = Presentation(uploaded_file)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"SLIDE {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            text = "\n".join(parts)

        elif file_type in IMAGE_EXT:
            img     = Image.open(uploaded_file)
            reader  = load_ocr()
            results = reader.readtext(img)
            text    = " ".join(r[1] for r in results)

    except Exception as e:
        st.error(f"Error reading file: {e}")
    return text.strip()

# =========================================================
# VECTORSTORE
# =========================================================
@st.cache_resource
def create_vectorstore(text: str):
    splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=300)
    chunks   = splitter.split_text(text)
    docs     = [LangDocument(page_content=c, metadata={"chunk_id": i})
                for i, c in enumerate(chunks)]
    emb      = HuggingFaceEmbeddings(model_name="BAAI/bge-large-en-v1.5")
    return FAISS.from_documents(docs, emb)

# =========================================================
# MODEL RUNNER — handles all pipeline task types uniformly
# =========================================================
def run_model(pipe_and_task, prompt: str) -> str:
    """Call the pipeline regardless of task type and return plain text."""
    pipe, task = pipe_and_task
    try:
        out = pipe(prompt, max_new_tokens=256)[0]
        if "generated_text" in out:
            text = out["generated_text"]
            # text-generation echoes the prompt — strip it
            if task == "text-generation" and text.startswith(prompt):
                text = text[len(prompt):].strip()
            return text.strip()
        if "summary_text" in out:
            return out["summary_text"].strip()
        return str(next(iter(out.values()))).strip()
    except Exception as e:
        return f"Model error: {e}"

# =========================================================
# RAG QUERY HELPER
# =========================================================
def rag_query(question, vectorstore, model, context_override=None):
    if context_override:
        context = context_override
        docs    = []
    elif vectorstore:
        docs    = vectorstore.similarity_search(question, k=5)
        context = "\n\n".join(d.page_content for d in docs)
    else:
        return "No index available.", []

    prompt = (
        "You are a precise document assistant. "
        "Answer ONLY from the context. "
        "If the answer is not in the context, say 'Not found in document'.\n\n"
        f"CONTEXT:\n{context}\n\nQUESTION: {question}\n\nANSWER:"
    )
    answer = run_model(model, prompt)
    return answer, docs

# =========================================================
# KEYWORD SEARCH across dataframe
# =========================================================
STOP = {"find","search","show","get","where","the","a","an","me","all","any","about",
        "for","with","in","of","and","or","that","which","what","list","tell","fetch",
        "display","filter","records","rows","give","entries","results","here"}

def df_search(question: str, df: pd.DataFrame) -> pd.DataFrame:
    q     = question.lower()
    terms = [w for w in re.split(r'\W+', q) if w not in STOP and len(w) > 2]
    matched = pd.DataFrame()
    for col in df.columns:
        for term in terms:
            temp    = df[df[col].astype(str).str.lower().str.contains(term, na=False)]
            matched = pd.concat([matched, temp])
    return matched.drop_duplicates().head(50)

# =========================================================
# UNIVERSAL ANSWER — one function, every file, every intent
# =========================================================
def universal_answer(question, file_type, text, df, vectorstore, model):
    """
    Returns (result_type, result_data)
    result_type : "markdown" | "table" | "dataframe_summary" | "rag"
    """
    intent = detect_intent(question)
    q      = question.lower().strip()

    # ── TABULAR (CSV / XLSX / XLS) ─────────────────────────
    if file_type in TABULAR and df is not None:

        if intent == "summary":
            return "dataframe_summary", df

        if intent == "columns":
            return "markdown", fmt_columns(df)

        if intent == "list":
            col = find_column(question, df)
            if col:
                return "markdown", fmt_list(df[col].dropna().tolist(), label=col)
            # no col found → list column names
            return "markdown", fmt_columns(df)

        if intent == "count":
            col = find_column(question, df)
            if col:
                n = df[col].nunique()
                return "markdown", fmt_count(f"{n} unique values", label=col)
            return "markdown", fmt_count(f"{df.shape[0]} rows, {df.shape[1]} columns")

        if intent == "table":
            results = df_search(question, df)
            if len(results):
                return "table", results
            return "markdown", "No matching records found."

        # factual — try column match first, then search
        col = find_column(question, df)
        if col:
            vals = df[col].dropna().tolist()
            # short answer if user asks for specific value
            if any(kw in q for kw in ["what is", "what's", "value of", "tell me"]):
                unique = sorted({str(v).strip() for v in vals if str(v).strip()})
                return "markdown", fmt_list(unique, label=col)
            return "markdown", fmt_list(vals, label=col)

        results = df_search(question, df)
        if len(results):
            return "table", results
        return "markdown", "No matching data found. Try mentioning a column name."

    # ── JSON ───────────────────────────────────────────────
    elif file_type == "json":
        try:
            data = json.loads(text)
        except Exception:
            return "markdown", "Could not parse JSON."

        if isinstance(data, dict):
            if intent == "list":
                return "markdown", fmt_list(list(data.keys()), label="JSON Keys")
            if intent == "count":
                return "markdown", fmt_count(len(data), label="Keys")
            keys_str = "\n".join(f"- {k}" for k in list(data.keys())[:30])
            return "markdown", f"**JSON — {len(data)} keys:**\n\n{keys_str}"

        elif isinstance(data, list):
            if intent == "count":
                return "markdown", fmt_count(len(data), label="Items")
            if intent == "list" and data and isinstance(data[0], dict):
                # try to find the field user wants
                tmp_df = pd.DataFrame(data[:1])
                col    = find_column(question, tmp_df)
                if col:
                    vals = [str(item.get(col,"")) for item in data if isinstance(item, dict)]
                    return "markdown", fmt_list(vals, label=col)
            return "markdown", f"JSON list — **{len(data)} items**.\n\nFirst item:\n```json\n{json.dumps(data[0], indent=2)[:600]}\n```"

    # ── IMAGE (OCR) ─────────────────────────────────────────
    elif file_type in IMAGE_EXT:
        if not text:
            return "markdown", "No text could be extracted from the image."
        if intent == "list":
            lines = [l.strip() for l in re.split(r'[\n,;]+', text) if l.strip() and len(l.strip()) > 1]
            if lines:
                return "markdown", fmt_list(lines, label="Extracted Items")
        if intent == "count":
            lines = [l.strip() for l in re.split(r'[\n,;]+', text) if l.strip()]
            return "markdown", fmt_count(len(lines), label="Extracted lines")
        return "markdown", f"**Extracted text from image:**\n\n{text}"

    # ── PPTX ───────────────────────────────────────────────
    elif file_type == "pptx":
        if intent == "list":
            lines = [l.strip() for l in text.split("\n")
                     if l.strip() and not l.strip().startswith("SLIDE") and len(l.strip()) > 3]
            return "markdown", fmt_list(lines, label="Slide Text Items")
        if intent == "summary":
            prompt = f"Summarize this PowerPoint presentation slide by slide:\n\n{text[:4000]}"
            result = run_model(model, prompt)
            return "markdown", result
        if intent == "count":
            slide_count = text.count("SLIDE ")
            return "markdown", fmt_count(slide_count, label="Slides")
        answer, docs = rag_query(question, vectorstore, model)
        return "rag", (answer, docs)

    # ── PDF / DOCX / TXT / MD / XML ─────────────────────────
    else:
        if intent == "list":
            # Pull lines from top RAG chunks that look like list items
            if vectorstore:
                rag_docs = vectorstore.similarity_search(question, k=6)
                context  = "\n".join(d.page_content for d in rag_docs)
                lines = [
                    l.strip() for l in context.split("\n")
                    if l.strip() and 3 < len(l.strip()) < 200
                ]
                if lines:
                    return "markdown", fmt_list(lines, label="Found Items")

        if intent == "count":
            answer, docs = rag_query(question, vectorstore, model)
            return "markdown", f"**Answer:** {answer}"

        if intent == "summary":
            prompt = f"Summarize this document clearly and concisely:\n\n{text[:4000]}"
            result = run_model(model, prompt)
            return "markdown", result

        answer, docs = rag_query(question, vectorstore, model)
        return "rag", (answer, docs)

    # Fallback
    return "markdown", "Could not determine an answer. Please rephrase your question."


# =========================================================
# MAIN UI
# =========================================================
st.title("🤖 Universal AI File Chatbot")
st.caption("Upload any file — ask any question — get the right answer in the right format.")

uploaded_file = st.file_uploader(
    "Upload a file",
    type=["pdf","docx","txt","csv","xlsx","xls","pptx","json","xml","md","png","jpg","jpeg"]
)

if uploaded_file:
    file_type = uploaded_file.name.split(".")[-1].lower()

    # Reset on new file
    if st.session_state.last_file != uploaded_file.name:
        st.session_state.update({**DEFAULTS, "last_file": uploaded_file.name})

    # Extract once
    if not st.session_state.text:
        with st.spinner("📖 Reading file..."):
            st.session_state.text      = extract_text(uploaded_file, file_type)
            st.session_state.file_type = file_type

    text      = st.session_state.text
    df        = st.session_state.df
    file_type = st.session_state.file_type

    if text:
        # Build vectorstore for non-tabular files
        if file_type not in TABULAR and st.session_state.vectorstore is None:
            with st.spinner("🔍 Building search index..."):
                st.session_state.vectorstore = create_vectorstore(text)

        model = load_model()

        # ── SIDEBAR ──────────────────────────────────────
        with st.sidebar:
            st.markdown("## 📁 File Info")
            st.success(f"✅ {uploaded_file.name}")
            st.write(f"**Type:** `{file_type.upper()}`")
            if df is not None:
                st.write(f"**Rows:** `{df.shape[0]}`")
                st.write(f"**Columns:** `{df.shape[1]}`")
                st.markdown("**Columns:**")
                for c in df.columns:
                    st.write(f"  • {c}")
            st.divider()
            st.markdown("**💡 Example questions:**")
            if file_type in TABULAR:
                st.caption("• List all company names\n• How many rows?\n• Find Pfizer\n• Show all emails\n• Summary")
            elif file_type == "pdf":
                st.caption("• Summarize the document\n• List all topics\n• What is the conclusion?")
            elif file_type in IMAGE_EXT:
                st.caption("• List all items\n• What text is in the image?")
            elif file_type == "pptx":
                st.caption("• How many slides?\n• Summarize the presentation\n• List all bullet points")
            elif file_type == "json":
                st.caption("• List all keys\n• How many items?\n• Show structure")

        # ── PREVIEW ──────────────────────────────────────
        st.subheader("📄 File Preview")
        if file_type in TABULAR and df is not None:
            st.dataframe(df.head(20), use_container_width=True, height=280)
        elif file_type in IMAGE_EXT:
            st.image(uploaded_file, use_column_width=True)
        else:
            st.text_area("Preview (first 3000 chars)", text[:3000], height=180)

        st.divider()

        # ── QUESTION ─────────────────────────────────────
        question = st.text_input(
            "💬 Ask anything about your file",
            placeholder="e.g.  List all company names  /  How many records?  /  Find Pfizer  /  Summary"
        )

        # ── ANSWER ───────────────────────────────────────
        if question:
            with st.spinner("🤔 Finding answer..."):
                result_type, result_data = universal_answer(
                    question, file_type, text, df,
                    st.session_state.vectorstore, model
                )

            st.subheader("🤖 Answer")

            # ── dataframe summary (metrics + table)
            if result_type == "dataframe_summary":
                c1, c2 = st.columns(2)
                c1.metric("Total Rows",    result_data.shape[0])
                c2.metric("Total Columns", result_data.shape[1])
                st.markdown("**Columns:**\n" + "\n".join(f"- {c}" for c in result_data.columns))
                st.dataframe(result_data.head(20), use_container_width=True, height=280)

            # ── filtered table (search results)
            elif result_type == "table":
                st.info(f"Found **{len(result_data)}** matching records")
                st.dataframe(result_data, use_container_width=True, height=350)

            # ── markdown (lists, counts, columns, summaries)
            elif result_type == "markdown":
                st.markdown(result_data)

            # ── rag (answer + optional source chunks)
            elif result_type == "rag":
                answer_text, rag_docs = result_data
                st.write(answer_text)
                if rag_docs:
                    with st.expander("📌 Source Chunks"):
                        for i, d in enumerate(rag_docs):
                            st.markdown(f"**Chunk {i+1}**")
                            st.write(d.page_content[:600])

    else:
        st.error("❌ Could not extract content from this file.")

else:
    st.markdown("""
    ### 👋 Welcome! Upload any of these file types:
    | Type | Extensions |
    |------|-----------|
    | 📊 Spreadsheet | CSV, XLSX, XLS |
    | 📄 Document | PDF, DOCX, TXT, MD |
    | 📑 Presentation | PPTX |
    | 🗂️ Data | JSON, XML |
    | 🖼️ Image | PNG, JPG, JPEG |
    """)