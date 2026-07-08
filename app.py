import streamlit as st
import tempfile
import os
import json
import numpy as np
import pandas as pd
import xml.etree.ElementTree as ET

from PIL import Image
from docx import Document
from pptx import Presentation
import easyocr

from langchain.schema import Document as LangDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter

from langchain_community.document_loaders import PyPDFLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

from transformers import pipeline


# =========================================================
# PAGE CONFIG
# =========================================================
st.set_page_config(
    page_title="Universal AI File Chatbot",
    page_icon="🤖",
    layout="wide"
)

# =========================================================
# SESSION STATE
# =========================================================
if "answer" not in st.session_state:
    st.session_state.answer = ""

if "docs" not in st.session_state:
    st.session_state.docs = []

if "last_file" not in st.session_state:
    st.session_state.last_file = None

# =========================================================
# CUSTOM CSS
# =========================================================
st.markdown("""
<style>

.block-container{
    padding-top:1rem;
    max-width:1100px;
}

h1{
    text-align:center;
}

.stButton button{
    border-radius:12px;
}

.stTextInput input{
    border-radius:12px;
}

</style>
""", unsafe_allow_html=True)

# =========================================================
# TITLE
# =========================================================
st.title("🤖 Universal AI File Chatbot")

st.write("""
Upload PDF, DOCX, TXT, CSV, XLSX, PPTX, JSON,
XML, Markdown or Images and ask questions using AI-powered RAG.
""")

# =========================================================
# OCR READER
# =========================================================
@st.cache_resource
def load_ocr():
    return easyocr.Reader(
        ["en"],
        gpu=False
    )

# =========================================================
# MODEL LOADER
# =========================================================
@st.cache_resource
def load_model():

    pipe = pipeline(
        "text-generation",
        model="gpt2",
        max_new_tokens=200
    )

    return pipe

# =========================================================
# HELPER: SAFE DATAFRAME DISPLAY
# =========================================================
def safe_df_display(df):
    """
    st.dataframe() uses a pyarrow backend under the hood.
    Columns with mixed types (e.g. dates mixed with strings/NaN,
    or Excel columns that come back as datetime.datetime objects
    inside an 'object' dtype column) crash pyarrow's type
    inference with errors like:
        ArrowTypeError: Expected bytes, got a 'datetime.datetime' object
    Converting everything to a clean string representation avoids
    this without needing to know the offending column ahead of time.
    """
    safe_df = df.copy()

    for col in safe_df.columns:
        try:
            # Try to keep numeric columns numeric for nicer display
            if pd.api.types.is_numeric_dtype(safe_df[col]):
                continue
            safe_df[col] = safe_df[col].astype(str)
        except Exception:
            safe_df[col] = safe_df[col].apply(
                lambda x: "" if pd.isna(x) else str(x)
            )

    return safe_df


def read_excel_any(file_obj, filename, sheet_name=None):
    """
    pd.read_excel needs the right engine depending on the file
    extension:
      - .xlsx / .xlsm -> openpyxl
      - .xls (legacy)  -> xlrd
    Without specifying this explicitly, legacy .xls uploads fail.
    """
    ext = filename.split(".")[-1].lower()
    engine = "xlrd" if ext == "xls" else "openpyxl"

    file_obj.seek(0)

    try:
        return pd.read_excel(file_obj, sheet_name=sheet_name, engine=engine)
    except Exception:
        # Fallback: let pandas guess if the explicit engine choice fails
        file_obj.seek(0)
        return pd.read_excel(file_obj, sheet_name=sheet_name)

# =========================================================
# TEXT EXTRACTION
# =========================================================
def extract_text(uploaded_file):

    file_type = uploaded_file.name.split(".")[-1].lower()

    text = ""

    try:

        # PDF
        if file_type == "pdf":

            uploaded_file.seek(0)

            with tempfile.NamedTemporaryFile(
                delete=False,
                suffix=".pdf"
            ) as tmp:

                tmp.write(uploaded_file.read())
                tmp_path = tmp.name

            loader = PyPDFLoader(tmp_path)

            docs = loader.load()

            text = "\n".join(
                [
                    doc.page_content.strip()
                    for doc in docs
                    if doc.page_content.strip()
                ]
            )

            os.unlink(tmp_path)

        # DOCX
        elif file_type == "docx":

            uploaded_file.seek(0)

            doc = Document(uploaded_file)

            text = "\n".join(
                [
                    para.text.strip()
                    for para in doc.paragraphs
                    if para.text.strip()
                ]
            )

        # TXT / MD
        elif file_type in ["txt", "md"]:

            uploaded_file.seek(0)

            text = uploaded_file.read().decode(
                "utf-8",
                errors="ignore"
            )

        # JSON
        elif file_type == "json":

            uploaded_file.seek(0)

            data = json.load(uploaded_file)

            text = json.dumps(
                data,
                indent=2
            )

        # XML
        elif file_type == "xml":

            uploaded_file.seek(0)

            tree = ET.parse(uploaded_file)

            root = tree.getroot()

            xml_text = []

            for elem in root.iter():

                if elem.text:

                    cleaned = elem.text.strip()

                    if cleaned:
                        xml_text.append(cleaned)

            text = "\n".join(xml_text)

        # CSV
        elif file_type == "csv":

            uploaded_file.seek(0)

            try:
                df = pd.read_csv(uploaded_file)
            except UnicodeDecodeError:
                uploaded_file.seek(0)
                df = pd.read_csv(uploaded_file, encoding="latin1")

            df = df.fillna("")

            rows = []

            for _, row in df.iterrows():

                row_text = " | ".join(
                    [
                        f"{col}: {str(row[col]).strip()}"
                        for col in df.columns
                        if str(row[col]).strip()
                    ]
                )

                rows.append(row_text)

            text = "\n".join(rows)

        # XLSX / XLS
        elif file_type in ["xlsx", "xls"]:

            excel_data = read_excel_any(
                uploaded_file,
                uploaded_file.name,
                sheet_name=None
            )

            all_rows = []

            for sheet_name, df in excel_data.items():

                df = df.fillna("")

                all_rows.append(
                    f"\n===== SHEET: {sheet_name} =====\n"
                )

                for _, row in df.iterrows():

                    row_text = " | ".join(
                        [
                            f"{col}: {str(row[col]).strip()}"
                            for col in df.columns
                            if str(row[col]).strip()
                        ]
                    )

                    if row_text:
                        all_rows.append(row_text)

            text = "\n".join(all_rows)

        # PPTX
        elif file_type == "pptx":

            uploaded_file.seek(0)

            prs = Presentation(uploaded_file)

            slides_text = []

            for i, slide in enumerate(prs.slides):

                slides_text.append(
                    f"===== SLIDE {i + 1} ====="
                )

                for shape in slide.shapes:

                    if hasattr(shape, "text"):

                        if shape.text.strip():

                            slides_text.append(
                                shape.text.strip()
                            )

            text = "\n".join(slides_text)

        # IMAGE OCR
        elif file_type in ["png", "jpg", "jpeg"]:

            uploaded_file.seek(0)

            image = Image.open(uploaded_file)
            image = image.convert("RGB")

            reader = load_ocr()

            # easyocr needs a NumPy array (or path/bytes) — NOT a PIL
            # Image object. Passing the PIL Image directly is the bug
            # that made png/jpg/jpeg extraction silently fail.
            image_np = np.array(image)

            results = reader.readtext(image_np)

            text = "\n".join(
                [
                    item[1]
                    for item in results
                ]
            )

        return text.strip()

    except Exception as e:

        st.error(
            f"❌ Error reading file: {str(e)}"
        )

        return ""

# =========================================================
# VECTORSTORE
# =========================================================
@st.cache_resource(show_spinner=False)
def create_vectorstore(text):

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100
    )

    chunks = splitter.split_text(text)

    chunks = [
        chunk.strip()
        for chunk in chunks
        if chunk.strip()
    ]

    docs = [
        LangDocument(page_content=chunk)
        for chunk in chunks
    ]

    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    vectorstore = FAISS.from_documents(
        docs,
        embeddings
    )

    return vectorstore

# =========================================================
# ASK QUESTION
# =========================================================
def ask_question(
    vectorstore,
    question,
    model
):

    retriever = vectorstore.as_retriever(
        search_kwargs={"k": 5}
    )

    docs = retriever.invoke(question)

    context = "\n\n".join(
        [
            doc.page_content
            for doc in docs
        ]
    )

    prompt = f"""
You are an intelligent AI document assistant.

Answer ONLY from the provided context.

Rules:
1. Give accurate answers.
2. Do not hallucinate.
3. If answer is not found say:
'I could not find that information in the uploaded file.'
4. Keep answers short and clear.

Context:
{context}

Question:
{question}

Answer:
"""

    try:

        result = model(prompt)

        answer = result[0]["generated_text"]

    except Exception as e:

        answer = f"Error generating answer: {str(e)}"

    return answer, docs
# =========================================================
# FILE UPLOADER
# =========================================================
uploaded_file = st.file_uploader(
    "📤 Upload File",
    type=[
        "pdf",
        "docx",
        "txt",
        "csv",
        "xlsx",
        "xls",
        "pptx",
        "json",
        "xml",
        "md",
        "png",
        "jpg",
        "jpeg"
    ]
)

# =========================================================
# MAIN APP
# =========================================================
if uploaded_file:

    if st.session_state.last_file != uploaded_file.name:

        st.session_state.answer = ""
        st.session_state.docs = []
        st.session_state.last_file = uploaded_file.name

    # =========================================
    # READ FILE
    # =========================================
    with st.spinner("📚 Reading file..."):
        text = extract_text(uploaded_file)

    if text:

        # =========================================
        # BUILD VECTORSTORE
        # =========================================
        with st.spinner("🧠 Building AI knowledge base..."):

            vectorstore = create_vectorstore(text)
            model = load_model()

        st.success("✅ File uploaded successfully!")

        file_type = uploaded_file.name.split(".")[-1].lower()

        # =========================================
        # FILE INFO
        # =========================================
        with st.sidebar:

            st.header("📁 File Information")

            st.write(
                f"**File Name:** {uploaded_file.name}"
            )

            st.write(
                f"**File Type:** {file_type.upper()}"
            )

            st.write(
                f"**Extracted Characters:** {len(text)}"
            )

        # =========================================
        # FILE PREVIEW
        # =========================================
        with st.expander("📄 File Preview", expanded=False):

            if file_type == "csv":

                uploaded_file.seek(0)

                try:
                    try:
                        df = pd.read_csv(uploaded_file)
                    except UnicodeDecodeError:
                        uploaded_file.seek(0)
                        df = pd.read_csv(uploaded_file, encoding="latin1")

                    st.dataframe(
                        safe_df_display(df),
                        use_container_width=True
                    )
                except Exception as e:
                    st.error(str(e))

            elif file_type in ["xlsx", "xls"]:

                try:

                    sheets = read_excel_any(
                        uploaded_file,
                        uploaded_file.name,
                        sheet_name=None
                    )

                    for sheet_name, df in sheets.items():

                        st.subheader(
                            f"📑 {sheet_name}"
                        )

                        st.dataframe(
                            safe_df_display(df),
                            use_container_width=True
                        )

                except Exception as e:
                    st.error(str(e))

            elif file_type in [
                "png",
                "jpg",
                "jpeg"
            ]:

                uploaded_file.seek(0)

                image = Image.open(uploaded_file)
                image = image.convert("RGB")

                st.image(
                    image,
                    use_container_width=True
                )

                st.markdown("### Extracted Text")

                st.write(text[:3000])

            elif file_type == "json":

                uploaded_file.seek(0)

                try:
                    data = json.load(uploaded_file)
                    st.json(data)
                except Exception:
                    st.write(text[:3000])

            else:

                st.write(text[:3000])

        # =========================================
        # QUESTION INPUT
        # =========================================
        st.markdown("---")

        question = st.text_input(
            "💬 Ask a question about the uploaded file"
        )

        # =========================================
        # GENERATE ANSWER
        # =========================================
        if question:

            with st.spinner("🔍 Searching document..."):

                answer, docs = ask_question(
                    vectorstore,
                    question,
                    model
                )

                st.session_state.answer = answer
                st.session_state.docs = docs

        # =========================================
        # DISPLAY ANSWER
        # =========================================
        if st.session_state.answer:

            st.subheader("🤖 Answer")

            st.write(
                st.session_state.answer
            )

            # =====================================
            # SOURCE CHUNKS
            # =====================================
            with st.expander(
                "📌 Source Chunks Used"
            ):

                for i, doc in enumerate(
                    st.session_state.docs
                ):

                    st.markdown(
                        f"### Chunk {i+1}"
                    )

                    st.write(
                        doc.page_content[:1000]
                    )

                    st.divider()

    else:

        st.error(
            "❌ No readable content found."
        )

else:

    st.info(
         "📂 Upload a file to begin."
    )